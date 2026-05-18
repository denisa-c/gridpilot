#!/usr/bin/env python3
"""
scripts/figures/make_architecture_pptx.py
=========================================
Produce two single-slide editable PowerPoint files:

  papers/whpc2026/architecture.pptx  --  GridPilot three-tier controller
  papers/pecs2026/architecture.pptx  --  f-SLA contract + AI baseline

Each .pptx contains one slide that the human author can open in
PowerPoint/Keynote/LibreOffice Impress to refine, then export to PDF
via File > Export > PDF.  The same scripts can be re-run to regenerate
the master copies from source.

The slides are intentionally simple --- coloured rectangles, plain
arrows, sans-serif labels --- so HPC experts who are not control or
mechanism-design specialists can grasp the architecture at a glance.

Dependencies
------------
  pip install python-pptx          # the only non-stdlib requirement

Usage
-----
  python scripts/figures/make_architecture_pptx.py
"""
from __future__ import annotations

import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:  # pragma: no cover
    print("ERROR: python-pptx is not installed.\n"
          "  pip install python-pptx\n"
          "Then re-run this script.")
    raise SystemExit(1)


# ─────────────────────────────────────────────────────────────────────
# Shared palette (print-safe; values chosen so a B&W print keeps the
# fill / border / arrow contrast distinguishable)
# ─────────────────────────────────────────────────────────────────────
COL_NAVY     = RGBColor(0x1E, 0x3A, 0x5F)   # deep navy header
COL_USER     = RGBColor(0xC8, 0xDC, 0xF0)   # light blue, fills user blocks
COL_OP       = RGBColor(0xFA, 0xE6, 0xC8)   # light gold, fills operator blocks
COL_LEDGER   = RGBColor(0xD8, 0xEB, 0xCD)   # mint, fills accounting blocks
COL_SCHED    = RGBColor(0xE5, 0xC5, 0xC0)   # dusty rose, fills scheduler block
COL_HW       = RGBColor(0xCF, 0xCF, 0xD3)   # neutral gray, fills hardware/grid blocks
COL_BORDER   = RGBColor(0x33, 0x33, 0x33)   # near-black borders
COL_LINE     = RGBColor(0x55, 0x55, 0x55)   # gray for arrows / dashed feedback
COL_TEXT     = RGBColor(0x1A, 0x1A, 0x1A)
COL_SUBTLE   = RGBColor(0x55, 0x55, 0x55)
COL_HILITE   = RGBColor(0xB8, 0x3A, 0x3A)   # red accent (callouts only)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────
# Small helpers
# ─────────────────────────────────────────────────────────────────────
def _new_presentation_169() -> Presentation:
    """16:9 presentation sized 13.33in x 7.5in (PowerPoint widescreen default)."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _add_box(slide, x, y, w, h, text, *, fill=COL_USER, border=COL_BORDER,
              font_size=12, bold=True, color=COL_TEXT, italic=False,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Add a rounded rectangle with centred text. Returns the shape."""
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(0.75)
    tf = shp.text_frame
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.text = ""
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(font_size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = "Calibri"
    return shp


def _add_label(slide, x, y, w, h, text, *, font_size=10,
                bold=False, italic=False, color=COL_TEXT,
                align=PP_ALIGN.CENTER):
    """Add a borderless text label (no shape outline)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.text = ""
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(font_size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


def _add_arrow(slide, x1, y1, x2, y2, *, color=COL_LINE, weight_pt=1.25,
                dashed=False, end_arrow=True):
    """Straight connector with optional dashed style and arrow head."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight_pt)
    if dashed:
        # python-pptx exposes prstDash via a low-level XML tweak
        from pptx.oxml.ns import qn
        ln = conn.line._get_or_add_ln()
        existing = ln.find(qn("a:prstDash"))
        if existing is not None:
            ln.remove(existing)
        dash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(dash)
    if end_arrow:
        from pptx.oxml.ns import qn
        ln = conn.line._get_or_add_ln()
        for tag in ("a:tailEnd", "a:headEnd"):
            existing = ln.find(qn(tag))
            if existing is not None:
                ln.remove(existing)
        tail = ln.makeelement(qn("a:tailEnd"),
                                {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
    return conn


def _add_title_bar(slide, title: str, subtitle: str | None = None):
    """Slide-wide header strip with bold title and optional subtitle."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0),
                                  SLIDE_W, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = COL_NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.4); tf.margin_right = Inches(0.4)
    tf.margin_top = Inches(0.10); tf.margin_bottom = Inches(0.0)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.text = ""
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run(); r1.text = title
    r1.font.size = Pt(26); r1.font.bold = True
    r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r1.font.name = "Calibri"
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(12); r2.font.italic = True
        r2.font.color.rgb = RGBColor(0xDD, 0xE6, 0xF2); r2.font.name = "Calibri"


def _add_footer(slide, text: str):
    """Small italic footer along the slide bottom."""
    _add_label(slide,
                Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3),
                text, font_size=9, italic=True, color=COL_SUBTLE,
                align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────────
# Slide 1: WHPC --- GridPilot three-tier controller
# ─────────────────────────────────────────────────────────────────────
def build_whpc_slide(prs: Presentation):
    """One slide: GridPilot three-tier controller, simplified."""
    blank = prs.slide_layouts[6]   # 'Blank' layout
    slide = prs.slides.add_slide(blank)

    _add_title_bar(
        slide,
        "GridPilot: Real-Time Grid-Responsive Control for AI Supercomputers",
        "Three control layers + an isolated 'safety island' that bypasses slow software for fast grid response."
    )

    # ── Layer labels (left margin)
    _add_label(slide, Inches(0.2), Inches(1.30), Inches(1.0), Inches(0.30),
                "Tier 3", font_size=11, bold=True, italic=True, color=COL_SUBTLE,
                align=PP_ALIGN.LEFT)
    _add_label(slide, Inches(0.2), Inches(2.55), Inches(1.0), Inches(0.30),
                "Tier 2", font_size=11, bold=True, italic=True, color=COL_SUBTLE,
                align=PP_ALIGN.LEFT)
    _add_label(slide, Inches(0.2), Inches(3.80), Inches(1.0), Inches(0.30),
                "Tier 1", font_size=11, bold=True, italic=True, color=COL_SUBTLE,
                align=PP_ALIGN.LEFT)
    _add_label(slide, Inches(0.2), Inches(5.30), Inches(1.0), Inches(0.30),
                "Hardware", font_size=11, bold=True, italic=True, color=COL_SUBTLE,
                align=PP_ALIGN.LEFT)

    # ── Tier 3 box (cluster operating point + PUE-aware FFR)
    _add_box(slide, Inches(1.3), Inches(1.25), Inches(7.5), Inches(1.0),
              "Cluster operating-point selector  (hourly)\n"
              "Picks a target IT power that meets the grid commitment\n"
              "at the meter --- includes the cooling-overhead correction.",
              fill=COL_SCHED, font_size=12)

    # ── Tier 2 box (host coordinator)
    _add_box(slide, Inches(1.3), Inches(2.55), Inches(7.5), Inches(1.0),
              "Per-host coordinator  (1 Hz)\n"
              "Predicts workload power one second ahead;\n"
              "splits the host envelope across GPUs.",
              fill=COL_LEDGER, font_size=12)

    # ── Tier 1 box (GPU PID)
    _add_box(slide, Inches(1.3), Inches(3.85), Inches(7.5), Inches(1.0),
              "Per-GPU power-cap loop  (200 Hz)\n"
              "Standard PID; tracks the assigned per-GPU power target\n"
              "via the NVIDIA Management Library.",
              fill=COL_USER, font_size=12)

    # ── Hardware box (GPUs + cooling)
    _add_box(slide, Inches(1.3), Inches(5.35), Inches(7.5), Inches(1.0),
              "GPU silicon + facility cooling\n"
              "Power swings settle on the GPU board within ~20 ms;\n"
              "facility meter reflects them within ~90 ms total.",
              fill=COL_HW, font_size=12)

    # ── Vertical arrows between layers (left column)
    for y1, y2 in [(2.25, 2.55), (3.55, 3.85), (4.85, 5.35)]:
        _add_arrow(slide, Inches(5.0), Inches(y1), Inches(5.0), Inches(y2))

    # ── Right column: grid + safety-island bypass
    _add_box(slide, Inches(9.3), Inches(1.25), Inches(3.6), Inches(1.0),
              "Electricity grid\nfrequency event",
              fill=COL_HW, font_size=13, color=COL_TEXT)

    _add_box(slide, Inches(9.3), Inches(3.40), Inches(3.6), Inches(1.4),
              "Safety island\n(real-time C bypass)\n"
              "Reads the grid trigger and pushes a new GPU power cap\n"
              "directly --- skipping the slower software path.",
              fill=COL_OP, font_size=11, italic=False)

    # Grid -> safety island
    _add_arrow(slide, Inches(11.1), Inches(2.25), Inches(11.1), Inches(3.40),
                weight_pt=1.5, color=COL_HILITE)
    # Safety island -> hardware (bypass)
    _add_arrow(slide, Inches(11.1), Inches(4.80), Inches(11.1), Inches(5.35),
                weight_pt=1.5, color=COL_HILITE)
    # Safety island -> hardware horizontal connector into the hardware box
    _add_arrow(slide, Inches(9.3), Inches(5.85), Inches(8.85), Inches(5.85),
                weight_pt=1.5, color=COL_HILITE)
    _add_label(slide, Inches(11.30), Inches(2.40), Inches(2.0), Inches(0.30),
                "trigger (UDP)", font_size=9, italic=True, color=COL_HILITE,
                align=PP_ALIGN.LEFT)
    _add_label(slide, Inches(11.30), Inches(4.85), Inches(2.0), Inches(0.30),
                "~ 97 ms median", font_size=9, italic=True, color=COL_HILITE,
                align=PP_ALIGN.LEFT)

    # ── Headline callout at the slide foot
    _add_box(slide, Inches(1.3), Inches(6.55), Inches(11.6), Inches(0.45),
              "Measured end-to-end response: 97 ms median (101 ms worst case across 90 trials)  ---  "
              "about 7x faster than the Nordic Fast Frequency Reserve 700 ms budget.",
              fill=RGBColor(0xF6, 0xEC, 0xEB), border=COL_HILITE,
              font_size=12, color=COL_HILITE, italic=True, bold=True,
              shape_type=MSO_SHAPE.RECTANGLE)

    _add_footer(slide,
        "Acronyms used in this slide:  PID = Proportional-Integral-Derivative controller  ·  "
        "GPU = Graphics Processing Unit  ·  PUE = Power Usage Effectiveness  ·  "
        "FFR = Fast Frequency Reserve.  "
        "Diagram is illustrative; the controller is open-source under MIT.")


# ─────────────────────────────────────────────────────────────────────
# Slide 2: PECS --- f-SLA + AI-baseline + gamification
# ─────────────────────────────────────────────────────────────────────
def build_pecs_slide(prs: Presentation):
    """One slide: f-SLA contract + AI baseline + gamification loop."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _add_title_bar(
        slide,
        "f-SLA: Eliciting User Flexibility for Carbon-Free Supercomputing",
        "An AI baseline predicts what the user will choose; the user picks a tier; the scheduler shifts the job to a clean hour."
    )

    # ── Top row: User game and Operator game
    _add_box(slide, Inches(0.5), Inches(1.20), Inches(5.7), Inches(1.40),
              "User submits a job\nand picks a tier on a 4-step ladder:\n"
              "T0 rigid · T1 hour · T2 day · T3 week",
              fill=COL_USER, font_size=13)

    _add_box(slide, Inches(7.1), Inches(1.20), Inches(5.7), Inches(1.40),
              "AI baseline (per-user predictor)\n"
              "shows the user the tier the AI expects;\n"
              "the user beats the AI to earn credit and rank.",
              fill=COL_OP, font_size=13)

    # ── Middle: f-SLA accounting layer
    _add_box(slide, Inches(0.5), Inches(2.95), Inches(12.3), Inches(1.10),
              "f-SLA accounting layer  ---  per-user credit ledger  +  leaderboard  +  log of (AI predicted, user declared, actually realised) for every job.\n"
              "This log is the experimental dataset: it replaces the synthetic tier prior used in this paper's proof-of-concept.",
              fill=COL_LEDGER, font_size=12)

    # ── Scheduler box
    _add_box(slide, Inches(0.5), Inches(4.40), Inches(12.3), Inches(1.20),
              "Carbon-aware scheduler  (any EASY/FCFS-class dispatcher)\n"
              "Defers each job within the user's declared tier window to a low-CI hour;\n"
              "the slowdown clause caps the worst-case wait.",
              fill=COL_SCHED, font_size=12)

    # ── Outcome box (electricity grid)
    _add_box(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.10),
              "Electricity grid (multi-country evaluation: SE, CH, FR, IT, DE, PL at 1 / 10 / 50 MW)\n"
              "CFE-lift ranges from +14 pp on a near-decarbonised grid (SE)\n"
              "down to +3.6 pp on a coal-heavy grid (PL).",
              fill=COL_HW, font_size=12)

    # ── Arrows
    # User-game -> accounting
    _add_arrow(slide, Inches(3.35), Inches(2.60), Inches(3.35), Inches(2.95))
    _add_label(slide, Inches(3.50), Inches(2.62), Inches(2.3), Inches(0.25),
                "declared tier", font_size=9, italic=True, color=COL_SUBTLE,
                align=PP_ALIGN.LEFT)
    # AI-baseline -> user-game (dashed, feedback)
    _add_arrow(slide, Inches(7.10), Inches(1.90), Inches(6.20), Inches(1.90),
                dashed=True)
    _add_label(slide, Inches(6.20), Inches(1.55), Inches(1.5), Inches(0.30),
                "AI prediction", font_size=9, italic=True, color=COL_SUBTLE,
                align=PP_ALIGN.LEFT)
    # AI-baseline -> accounting (logs triples)
    _add_arrow(slide, Inches(9.95), Inches(2.60), Inches(9.95), Inches(2.95))
    # accounting -> scheduler
    _add_arrow(slide, Inches(6.65), Inches(4.05), Inches(6.65), Inches(4.40))
    # scheduler -> grid
    _add_arrow(slide, Inches(6.65), Inches(5.60), Inches(6.65), Inches(5.95))

    # ── Bottom right: PoC stamp
    _add_label(slide, Inches(7.5), Inches(7.10), Inches(5.4), Inches(0.30),
                "Proof of concept --- numbers are illustrative; the contract is the contribution.",
                font_size=9, italic=True, color=COL_HILITE,
                align=PP_ALIGN.RIGHT)
    _add_footer(slide,
        "Acronyms used in this slide:  f-SLA = flexible Service-Level Agreement  ·  "
        "CFE = Carbon-Free Energy  ·  CI = Carbon Intensity  ·  "
        "EASY/FCFS = standard HPC queueing policies.")


# ─────────────────────────────────────────────────────────────────────
# Entry point
# ─────────────────────────────────────────────────────────────────────
def main(argv=None):
    parser = argparse.ArgumentParser(description=__doc__,
                                     formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("--papers-root", type=Path,
                        default=Path(__file__).resolve().parents[3] / "papers",
                        help="Root of the papers/ directory (default: ../../papers).")
    args = parser.parse_args(argv)

    whpc_out = args.papers_root / "whpc2026" / "architecture.pptx"
    pecs_out = args.papers_root / "pecs2026" / "architecture.pptx"
    whpc_out.parent.mkdir(parents=True, exist_ok=True)
    pecs_out.parent.mkdir(parents=True, exist_ok=True)

    prs_w = _new_presentation_169(); build_whpc_slide(prs_w); prs_w.save(str(whpc_out))
    print(f"[arch-pptx] wrote {whpc_out}")

    prs_p = _new_presentation_169(); build_pecs_slide(prs_p); prs_p.save(str(pecs_out))
    print(f"[arch-pptx] wrote {pecs_out}")

    print("\nNext: open in PowerPoint/Keynote/LibreOffice, refine, then\n"
          "  File > Export > PDF  ->  papers/<paper>/figs/architecture.pdf\n"
          "The papers' build.sh stages architecture.pdf into the LaTeX figs/ dir.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
